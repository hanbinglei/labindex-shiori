"""
LabIndex Shiori — PPTX 解析器

功能:
  1. 抽取标题（PPT 内置属性优先 → 首页文本框推测兜底）
  2. 抽取关键词（キーワード:/Keywords: 标注优先 → 频率推测兜底）
  3. 抽取正文文本（各幻灯片文本框，供全文检索）
  4. 抽取年份、作者
  5. 来源标注: annotated / inferred

只读保障: 仅 open(path, 'rb')
fail-safe: 损坏/加密/异常 PPT → 记日志跳过不中断
"""

from __future__ import annotations

import logging
import re
from collections import Counter
from typing import Dict, List, Optional, Tuple

# 复用 PDF 解析器的关键词/年份/停用词工具
from src.parser.pdf_parser import (
    _KEYWORD_PATTERNS,
    _STOPWORDS,
    _extract_keywords,
    _extract_year_from_text,
    _parse_authors,
)

logger = logging.getLogger(__name__)

# PPT 标题应过滤的干扰模式（页眉、页码、日期等）
_TITLE_SKIP_PATTERNS = [
    r"^\d+$",                           # 纯数字
    r"^\d{4}[/\-]\d{2}[/\-]\d{2}",      # 日期 2024/01/15
    r"^第[\d]+章",                       # 第X章
    r"^[A-Z\s]{8,}$",                   # 全大写短句（可能是机构名）
    r"^目次$", r"^目錄$", r"^目录$",     # 目录页
    r"^概要$", r"^Abstract$",            # 章节标题
    r"^参考文献$", r"^References$",       # 章末标题
    r"^ご清聴", r"^Thank you",           # 结尾页
    r"^付録", r"^Appendix",              # 附录
]

# 日文/英文关键词标注模式（与 PDF 相同，增加 PPT-specific 格式）
_PPT_KEYWORD_PATTERNS = _KEYWORD_PATTERNS + [
    r"キーワード[：:]\s*(.+?)(?:\n|$)",
    r"Keywords?[：:]\s*(.+?)(?:\n|$)",
]


def parse_pptx(file_path: str, max_slides: int = 200) -> Dict:
    """
    解析 PPTX 文件，返回结构化元数据

    Args:
        file_path: PPTX 文件路径
        max_slides: 最大解析幻灯片数（防超大文件卡死）

    Returns:
        dict: {
            "title": str or None,
            "title_source": "annotated" or "inferred",
            "keywords": list[str] or None,
            "keywords_source": "annotated" or "inferred",
            "abstract": str or None,       # PPT 的正文全文
            "year": int or None,
            "authors": list[str] or None,
            "references": None,            # PPT 通常无参考文献章节
            "note": str or None,
            "extra": dict,
        }
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    total_slides = min(len(prs.slides), max_slides)

    # --- 获取 PPT 内置属性 ---
    core_props = prs.core_properties
    ppt_title = getattr(core_props, "title", "") or ""
    ppt_author = getattr(core_props, "author", "") or ""

    # --- 逐幻灯片抽取文本 ---
    all_text_parts = []     # 全部文本（供全文检索/关键词抽取）
    first_slide_text = ""   # 首页文本（供标题推测）

    for idx, slide in enumerate(prs.slides):
        if idx >= max_slides:
            logger.info("PPTX 达到幻灯片上限 %d，截断 [%s]", max_slides, file_path)
            break

        slide_texts = []
        try:
            for shape in slide.shapes:
                # 只处理有文本的 shape
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                # 也尝试从 table 中提取文本
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_texts.append(cell.text.strip())
        except Exception as e:
            # 单个 slide 异常不中断整体
            logger.warning("PPTX slide %d 处理异常 [%s]: %s", idx, file_path, e)
            continue

        slide_full = "\n".join(slide_texts)
        all_text_parts.append(slide_full)

        # 记录首页文本
        if idx == 0:
            first_slide_text = slide_full

    full_text = "\n".join(all_text_parts)

    # --- 检查是否有可提取的文本 ---
    if not full_text.strip():
        return {
            "title": ppt_title.strip() or None,
            "title_source": "annotated" if ppt_title.strip() else "inferred",
            "keywords": None,
            "keywords_source": None,
            "abstract": None,
            "year": None,
            "authors": _parse_authors(ppt_author) if ppt_author else None,
            "references": None,
            "note": "PPTX 文件无可提取文本",
            "extra": {},
        }

    # --- 1. 标题 ---
    title, title_source = _extract_pptx_title(first_slide_text, ppt_title)

    # --- 2. 关键词 ---
    keywords, kw_source = _extract_pptx_keywords(full_text)

    # --- 3. 正文（供全文检索）---
    # PPT 没有传统意义的摘要；将全部文本作为 abstract 字段存入
    abstract = full_text[:5000] if full_text.strip() else None

    # --- 4. 年份 ---
    year = _extract_year_from_text(full_text)

    # --- 5. 作者 ---
    authors = _parse_authors(ppt_author) if ppt_author else None

    return {
        "title": title,
        "title_source": title_source,
        "keywords": keywords,
        "keywords_source": kw_source,
        "abstract": abstract,
        "year": year,
        "authors": authors,
        "references": None,  # PPT 通常无参考文献章节
        "note": None,
        "extra": {},
    }


# ============================================================
# 标题提取
# ============================================================

def _extract_pptx_title(first_slide_text: str, ppt_title: str) -> Tuple[Optional[str], str]:
    """
    从 PPT 内置属性或首页文本提取标题

    策略:
    1. PPT 内置属性 title（非空且合理）→ annotated
    2. 首页第一段有意义的文本框内容 → inferred
    3. 过滤干扰模式（日期、页码、目录等）

    Returns:
        (title_text, source_label)
    """
    # 方案1: PPT 内置标题
    if ppt_title and len(ppt_title.strip()) > 5:
        if not re.match(r"^(untitled|无题|PPT|Presentation|PowerPoint|default)", ppt_title, re.I):
            return ppt_title.strip(), "annotated"

    # 方案2: 从首页文本框提取
    if not first_slide_text.strip():
        return None, "inferred"

    lines = first_slide_text.strip().split("\n")
    lines = [l.strip() for l in lines if l.strip()]

    for line in lines:
        if len(line) < 5:
            continue
        # 跳过干扰模式
        if any(re.search(p, line) for p in _TITLE_SKIP_PATTERNS):
            continue
        return line, "inferred"

    return None, "inferred"


# ============================================================
# 关键词提取
# ============================================================

def _extract_pptx_keywords(full_text: str) -> Tuple[Optional[List[str]], Optional[str]]:
    """
    从 PPT 全文提取关键词

    优先级:
    1. annotated: 找到 キーワード:/Keywords: 标注
    2. inferred: 无标注，TF 频率提取特征词

    复用了 PDF 的 _extract_keywords 逻辑，但增加了 PPT-specific 模式。
    """
    # 尝试标注提取（用扩展后的模式）
    for pattern in _PPT_KEYWORD_PATTERNS:
        m = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)
        if m:
            kw_text = m.group(1).strip()
            kws = re.split(r"[,、，・\s/／]+", kw_text)
            kws = [k.strip().strip(".") for k in kws if k.strip() and len(k.strip()) > 1]
            kws = [k for k in kws if not re.match(r"^\d+$", k)]
            if kws:
                return kws, "annotated"

    # 兜底：复用 PDF 的算法推测
    return _extract_keywords(full_text)
